from flask import Flask, render_template, request
import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from transformers import pipeline
import random

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'

# Initialize summarizer
summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

# Extract text from different file types
def extract_text(file_path, file_type):
    text = ""
    if file_type == 'pdf':
        reader = PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    elif file_type == 'docx':
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file_type == 'pptx':
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

# Generate flashcards (simple split into Q/A)
def generate_flashcards(text):
    lines = text.split('.')
    flashcards = []
    for line in lines:
        if len(line.strip()) > 20:
            question = line.strip()[:50] + "?"
            answer = line.strip()
            flashcards.append((question, answer))
    return flashcards

# Generate simple quiz (fill-in-the-blank MCQs)
def generate_quiz(text, num_questions=5):
    sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 20]
    quiz = []
    if len(sentences) < num_questions:
        num_questions = len(sentences)
    for _ in range(num_questions):
        sentence = random.choice(sentences)
        words = sentence.split()
        if len(words) < 5:
            continue
        answer = random.choice(words)
        question = sentence.replace(answer, "_____")
        # Generate 3 random options
        options = [answer]
        while len(options) < 4:
            word = random.choice(words)
            if word not in options:
                options.append(word)
        random.shuffle(options)
        quiz.append({"question": question, "answer": answer, "options": options})
    return quiz

@app.route('/', methods=['GET', 'POST'])
def index():
    notes = ""
    flashcards = []
    quiz = []
    if request.method == 'POST':
        file = request.files['file']
        if file:
            filename = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(filename)

            # Detect file type
            ext = file.filename.split('.')[-1].lower()
            text = extract_text(filename, ext)

            # Summarize notes
            summarized = summarizer(text, max_length = 10000, min_length=50, do_sample=False)
            notes = summarized[0]['summary_text']

            # Generate flashcards
            flashcards = generate_flashcards(notes)

            # Generate quiz
            quiz = generate_quiz(notes)

    return render_template('index.html', notes=notes, flashcards=flashcards, quiz=quiz)

if __name__ == '__main__':
    app.run(debug=True)
